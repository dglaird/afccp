from pptx import Presentation
import io
import matplotlib.pyplot as plt
from pptx.util import Inches, Pt
from PIL import Image
import os
import numpy as np

# afccp modules
import afccp.globals


# __________________________________________________RESULTS_____________________________________________________________
def generate_results_slides(instance):
    """
    Function to generate the results slides for a particular problem instance with solution
    """

    # Shorthand
    mdl_p = instance.mdl_p

    # Build the presentation object
    prs = Presentation(afccp.globals.paths['files'] + 'results_slide_template.pptx')

    # Delete the current slides in the template
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    # Slide Layouts
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]
    chart_slide_layout = prs.slide_layouts[3]
    bubble_slide_layout = prs.slide_layouts[4]
    closing_slide_layout = prs.slide_layouts[5]

    # Set width and height of presentation
    prs.slide_width = Inches(mdl_p['b_figsize'][0])
    prs.slide_height = Inches(mdl_p['b_figsize'][1])

    # Add title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = instance.data_name + " Classification Results (" + instance.solution['name'] + ")"
    # print(len(slide.placeholders))
    # print(slide.placeholders)
    # content = slide.placeholders[1]
    # content.text = "Rank, Name\nAFPC/DSYA\nTBD"

    # Initial content slide
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    title.text = "Overview"
    content = slide.placeholders[1]
    content.text = "Here's where the overview goes"

    # Size of the chart
    # top, left = Inches(instance.mdl_p['ch_top']), Inches(instance.mdl_p['ch_left'])
    # height, width = Inches(instance.mdl_p['ch_height']), Inches(instance.mdl_p['ch_width'])

    # # Size of the chart
    # top, left = Inches(0), Inches(0)
    # height, width = Inches(mdl_p['figsize'][1]), Inches(mdl_p['figsize'][0])

    # Get the file paths to all the relevant images
    folder_path = instance.export_paths['Analysis & Results'] + instance.solution_name + "/"
    folder = os.listdir(folder_path)
    chart_paths = {}
    for file in folder:
        if instance.solution_name in file and '.png' in file:
            chart_paths[file] = folder_path + file

    # Loop through the pictures I want in the order I want
    chart_text_order = {"PGL": ["Combined Quota", "quantity_bar"], "BUBBLE CHART": ["Cadet Choice.png"],
                        "SOC by Accessions Group": ["Accessions", "SOC Chart"],
                        "Gender by Accessions Group": ["Accessions", "Gender Chart"],
                        "Race by Accessions Group": ["Accessions", "Race Chart"],
                        "Ethnicity by Accessions Group": ["Accessions", "Ethnicity Chart"],
                        "SOC by AFSC": ["Extra Measure", "SOC Chart_proportion"],
                        "Gender by AFSC": ["Extra Measure", "Gender Chart_proportion"],
                        "Race by AFSC": ["Extra Measure", "Race Chart_proportion"],
                        "Ethnicity by AFSC": ["Extra Measure", "Ethnicity Chart_proportion"],
                        "Cadet Preference by AFSC": ["Utility", "quantity_bar_choice"],
                        "AFSC Preference": ["Norm Score", "quantity_bar_choice"],
                        "62EXE BUBBLE CHART": ["62EXE Specific Choice.png"],}
    for title_text in chart_text_order:

        # Loop through each potential image
        for file in chart_paths:

            # Determine if this is the image I want here
            found = True
            for test_text in chart_text_order[title_text]:
                if test_text not in file:
                    found = False
                    break

            # If this is the file I want, we do things
            if found:

                # Determine the layout of the chart needed
                if "BUBBLE CHART" in title_text:

                    # Create the bubble chart slide
                    slide = prs.slides.add_slide(bubble_slide_layout)

                    # Add the image to the slide
                    for shape in slide.placeholders:
                        if "Picture" in shape.name:
                            shape.insert_picture(chart_paths[file])
                else:

                    # Create the AFSC/Accessions Chart slide
                    slide = prs.slides.add_slide(chart_slide_layout)
                    title = slide.shapes.title
                    title.text = title_text

                    # # Add the picture to the slide
                    # slide.shapes.add_picture(chart_paths[file], left, top, height=height, width=width)

                    # Add the image to the slide
                    for shape in slide.placeholders:
                        if "Picture" in shape.name:
                            shape.insert_picture(chart_paths[file])

                # Break out of this loop since we found the image we want
                break


    # Add closing slide
    prs.slides.add_slide(closing_slide_layout)

    # Save the PowerPoint
    filepath = instance.export_paths['Analysis & Results'] + instance.solution_name + '/' + \
               instance.data_name + ' ' + instance.solution_name + '.pptx'
    prs.save(filepath)


def generate_comparison_slides(instance):
    """
    Function to generate the results slides for a particular problem instance with solution
    """

    # Shorthand
    mdl_p = instance.mdl_p

    # Build the presentation object
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    # Set width and height of presentation
    prs.slide_width = Inches(mdl_p['b_figsize'][0])
    prs.slide_height = Inches(mdl_p['b_figsize'][1])

    # Size of the chart
    top, left = Inches(0), Inches(0)

    # Get the file paths to all the relevant images
    folder_path = instance.export_paths['Analysis & Results'] + "Comparison Charts/"
    folder = os.listdir(folder_path)

    # Sort the files in the folder according to my preferred method
    chart_paths = {}
    keyword_order = [['Combined Quota'], ['Tier 1'], ['Male'], ['USAFA Proportion'], ['Merit'], ['Norm Score'],
                     ['Utility', 'dot'], ['Utility', 'mean_preference'], ['Utility', 'median_preference'],
                     ['Utility', 'Histogram'], ['Pareto', 'Cadets.png'], ['Pareto', ').png'], ['Similarity']]

    # Loop through each set of "keywords"
    for keywords in keyword_order:

        # Loop through each file until we have a match
        for file in folder:

            # Loop through all keywords to see if we have a match
            match = True
            for keyword in keywords:
                if keyword not in file:
                    match = False
                    break

            # If we have a "match", add it!
            if match:
                chart_paths[file] = folder_path + file
                folder.remove(file)  # Remove this file since we've accounted for it
                break

    # Add any remaining files
    for file in folder:
        if '.png' in file:
            chart_paths[file] = folder_path + file

    # Loop through each image file path to add the image to the presentation
    for file in chart_paths:

        # Add an empty slide
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add the picture to the slide
        slide.shapes.add_picture(chart_paths[file], left, top)  #, height=height, width=width)

    # Save the PowerPoint
    filepath = instance.export_paths['Analysis & Results'] + instance.data_name + ' ' + 'Comparison.pptx'
    prs.save(filepath)


def create_animation_slides(instance):
    """
    Function to generate the results slides for a particular problem instance with solution
    """

    # Shorthand
    mdl_p, sequence = instance.mdl_p, instance.solution['iterations']['sequence']

    # Build the presentation object
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[7]

    # Set width and height of presentation
    prs.slide_width = Inches(mdl_p['b_figsize'][0])
    prs.slide_height = Inches(mdl_p['b_figsize'][1])

    # Size of the chart
    top, left = Inches(0), Inches(0)
    height, width = Inches(mdl_p['b_figsize'][1]), Inches(mdl_p['b_figsize'][0])

    # Make sure we have the "sequence" folder
    if sequence not in os.listdir(instance.export_paths['Analysis & Results'] + "Cadet Board/"):
        raise ValueError("Error. Sequence folder '" + sequence + "' not in 'Cadet Board' folder.")

    # Get the file path to the sequence folder
    sequence_folder_path = instance.export_paths['Analysis & Results'] + "Cadet Board/" + sequence + '/'

    # Make sure we have the "focus" folder
    if mdl_p['focus'] not in os.listdir(sequence_folder_path):
        raise ValueError("Error. Focus folder '" + mdl_p['focus'] + "' not in '" + sequence + "' sequence folder.")

    # Files in the focus folder
    focus_folder_path = sequence_folder_path + mdl_p['focus'] + '/'
    folder = np.array(os.listdir(focus_folder_path))

    # Regular solution iterations frames: 1.png for example
    if "1.png" in folder:

        # Sort the folder in order by the frames
        int_vals = np.array([int(file[0]) for file in folder])
        indices = np.argsort(int_vals)
        img_paths = {file: focus_folder_path + file for file in folder[indices]}

    # Matching Algorithm proposals/rejections frames: 1 (Proposals).png & 1 (Rejections).png for example
    else:

        # The integer values at the beginning of each file
        int_vals = np.array([int(file[:2]) for file in folder])
        min, max = np.min(int_vals), np.max(int_vals)

        # Loop through the files to get the ordered list of frames
        img_paths = {}
        for val in np.arange(min, max + 1):
            indices = np.where(int_vals == val)[0]
            files = folder[indices]
            if len(files) > 1:
                if 'Proposals' in files[0]:
                    img_paths[files[0]] = focus_folder_path + files[0]
                    img_paths[files[1]] = focus_folder_path + files[1]
                else:
                    img_paths[files[1]] = focus_folder_path + files[1]
                    img_paths[files[0]] = focus_folder_path + files[0]
            else:
                img_paths[files[0]] = focus_folder_path + files[0]

    # Loop through each image file path to add the image to the presentation
    for file in img_paths:

        # Add an empty slide
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add the picture to the slide
        slide.shapes.add_picture(img_paths[file], left, top, height=height, width=width)

    # Save the PowerPoint
    filepath = sequence_folder_path + mdl_p['focus'] + '.pptx'
    prs.save(filepath)


def create_animated_presentation(instance, num_intro_slides=3):
    """
    Generates a PowerPoint presentation with:
    - Orientation slide (if available)
    - First few still frames
    - Animated GIF slide with remaining frames
    - Final solution slide (if available)

    Parameters:
        instance: Object containing model data and paths
        num_intro_slides (int): Number of still slides to include before GIF
    """
    # Shorthand
    mdl_p = instance.mdl_p
    sequence = instance.solution['iterations']['sequence']

    # Build the presentation
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[7]
    prs.slide_width = Inches(mdl_p['b_figsize'][0])
    prs.slide_height = Inches(mdl_p['b_figsize'][1])
    top, left = Inches(0), Inches(0)
    height, width = Inches(mdl_p['b_figsize'][1]), Inches(mdl_p['b_figsize'][0])

    # Paths
    base_path = instance.export_paths['Analysis & Results'] + "Cadet Board/"
    sequence_path = os.path.join(base_path, sequence)
    focus_path = os.path.join(sequence_path, mdl_p['focus'])

    # Get sorted image files
    folder = np.array(os.listdir(focus_path))
    image_files = sorted([f for f in folder if f.endswith('.png')],
                         key=lambda x: int(''.join(filter(str.isdigit, x.split()[0]))))

    # --- Orientation Slide ---
    orientation_path = os.path.join(focus_path, 'orientation.png')
    if os.path.exists(orientation_path):
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(orientation_path, left, top, height=height, width=width)

    # --- Still Slides ---
    for i, file in enumerate(image_files[:num_intro_slides]):
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(os.path.join(focus_path, file), left, top, height=height, width=width)

    # --- Create GIF from remaining images ---
    gif_images = []
    for file in image_files[num_intro_slides:]:
        img = Image.open(os.path.join(focus_path, file)).convert("RGB")
        gif_images.append(img)

    if gif_images:
        gif_path = os.path.join(focus_path, 'animation.gif')
        gif_images[0].save(
            gif_path,
            save_all=True,
            append_images=gif_images[1:],
            duration=500,
            loop=0
        )

        # Add GIF slide (note: PowerPoint may require manual replacement depending on viewer support)
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(gif_path, left, top, height=height, width=width)

    # --- Final Image (if exists) ---
    file = image_files[len(image_files) - 1]
    slide = prs.slides.add_slide(blank_slide_layout)
    slide.shapes.add_picture(os.path.join(focus_path, file), left, top, height=height, width=width)

    # Save presentation
    output_path = os.path.join(sequence_path, mdl_p['focus'] + '_animated.pptx')
    prs.save(output_path)


# __________________________________________________CADET PROFILES______________________________________________________
def generate_cadet_profile_slides(instance):

    # Shorthand
    mdl_p = instance.mdl_p
    p = instance.parameters
    cadet = mdl_p['cadet']

    # Build the presentation object
    prs = Presentation(afccp.globals.paths['files'] + 'results_slide_template.pptx')

    # Delete the current slides in the template
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    # Slide Layouts
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]
    two_column_slide_layout = prs.slide_layouts[2]
    chart_slide_layout = prs.slide_layouts[3]
    bubble_slide_layout = prs.slide_layouts[4]
    closing_slide_layout = prs.slide_layouts[5]

    # Set width and height of presentation
    prs.slide_width = Inches(mdl_p['b_figsize'][0])
    prs.slide_height = Inches(mdl_p['b_figsize'][1])

    # Add title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = f"{instance.data_name} Cadet '{cadet}' One Market Model Profile"

    # Initial content slide
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    title.text = "Overview"
    content = slide.placeholders[1]
    content.text = "Here's where the overview goes"

    # Add the cadet preference submission slide
    _add_cadet_afsc_base_course_preferences_overview_slide(instance, cadet, prs)

    # Add the cadet's utility chart
    if 'J^Selected' in p:  # This is here cause some data instances can't support this slide
        _add_cadet_afsc_utility_chart_slide(instance, cadet, prs)
    if 'bases' in p:
        _add_cadet_state_utility_chart_slide(instance, cadet, prs)
        _add_base_utility_chart_slide(instance, cadet, prs)
        _add_course_utility_chart_slide(instance, cadet, prs)

    # Add closing slide
    prs.slides.add_slide(closing_slide_layout)

    # Save the powerpoint
    folder_path = instance.export_paths['Analysis & Results'] + 'Data Charts' + '/'
    if 'Cadet Profile Slides' not in os.listdir(folder_path):
        os.mkdir(folder_path + 'Cadet Profile Slides')
    prs.save(f'{folder_path}Cadet Profile Slides/{cadet}.pptx')


def _add_cadet_afsc_base_course_preferences_overview_slide(instance, cadet, prs):

    from pptx.util import Pt
    import numpy as np

    p = instance.parameters
    i = cadet

    slide = prs.slides.add_slide(prs.slide_layouts[2])  # your custom layout

    # ----------------------------------------------------------
    # Placeholders (based on your layout)
    # ----------------------------------------------------------

    title_ph   = slide.placeholders[0]
    left_ph    = slide.placeholders[1]
    right_ph   = slide.placeholders[12]
    summary_ph = slide.placeholders[13]

    title_ph.text = f"Cadet '{i}' — Preference & Decision Profile"

    # ----------------------------------------------------------
    # Compute Metrics
    # ----------------------------------------------------------

    selected = p['J^Selected'][i]
    eligible = p['J^E'][i]

    eligible_selected = np.intersect1d(selected, eligible)
    ineligible_selected = np.setdiff1d(selected, eligible)

    selected_count = len(selected)
    eligible_count = len(eligible_selected)
    eligible_pct = round(100 * eligible_count / selected_count, 1) if selected_count > 0 else 0

    bottom_choices = p['J^Bottom 2 Choices'][i]
    last_choice = p['J^Last Choice'][i]

    bottom_names = [p['afscs'][j] for j in bottom_choices]
    last_name = p['afscs'][last_choice]
    ineligible_names = [p['afscs'][j] for j in ineligible_selected]

    bt = p['base_threshold'][i]
    tt = p['training_threshold'][i]
    training_pref = p['training_preferences'][i]

    wA = round(p['weight_afsc'][i], 2)
    wB = round(p['weight_base'][i], 2)
    wC = round(p['weight_course'][i], 2)

    # ----------------------------------------------------------
    # LEFT COLUMN (Larger Font)
    # ----------------------------------------------------------

    left_tf = left_ph.text_frame
    left_tf.clear()

    p0 = left_tf.paragraphs[0]
    p0.text = "AFSC Submission Overview"
    p0.font.bold = True
    p0.font.size = Pt(18)

    body_font_size = Pt(15)

    p1 = left_tf.add_paragraph()
    p1.text = f"AFSCs Submitted: {selected_count}"
    p1.font.size = body_font_size

    p2 = left_tf.add_paragraph()
    p2.text = f"Eligible Among Submitted: {eligible_count} ({eligible_pct}%)"
    p2.level = 1
    p2.font.size = body_font_size

    if len(ineligible_names) > 0:
        p3 = left_tf.add_paragraph()
        p3.text = "Selected but Not Eligible:"
        p3.font.bold = True
        p3.font.size = body_font_size

        for name in ineligible_names:
            pb = left_tf.add_paragraph()
            pb.text = name
            pb.level = 1
            pb.font.size = body_font_size

    p4 = left_tf.add_paragraph()
    p4.text = "Bottom Ranked AFSCs:"
    p4.font.bold = True
    p4.font.size = body_font_size

    for name in bottom_names:
        pb = left_tf.add_paragraph()
        pb.text = name
        pb.level = 1
        pb.font.size = body_font_size

    p5 = left_tf.add_paragraph()
    p5.text = f"Final Ranked AFSC: {last_name}"
    p5.font.size = body_font_size

    # ----------------------------------------------------------
    # RIGHT COLUMN (Larger Font)
    # ----------------------------------------------------------

    right_tf = right_ph.text_frame
    right_tf.clear()

    r0 = right_tf.paragraphs[0]
    r0.text = "Decision Parameters"
    r0.font.bold = True
    r0.font.size = Pt(18)

    r1 = right_tf.add_paragraph()
    r1.text = f"Training Preference: {training_pref}"
    r1.font.size = body_font_size

    r2 = right_tf.add_paragraph()
    r2.text = "Threshold Settings"
    r2.font.bold = True
    r2.font.size = body_font_size

    r3 = right_tf.add_paragraph()
    r3.text = f"Base Threshold: {bt}"
    r3.level = 1
    r3.font.size = body_font_size

    r4 = right_tf.add_paragraph()
    r4.text = f"Training Threshold: {tt}"
    r4.level = 1
    r4.font.size = body_font_size

    r5 = right_tf.add_paragraph()
    r5.text = "Objective Weights"
    r5.font.bold = True
    r5.font.size = body_font_size

    for text in [
        f"AFSC Weight: {wA}",
        f"Base Weight: {wB}",
        f"Course Weight: {wC}",
    ]:
        pr = right_tf.add_paragraph()
        pr.text = text
        pr.level = 1
        pr.font.size = body_font_size

    # ----------------------------------------------------------
    # SUMMARY (No Bullet)
    # ----------------------------------------------------------

    summary_text = _generate_cadet_summary_text(p, i)

    summary_tf = summary_ph.text_frame
    summary_tf.clear()

    s0 = summary_tf.paragraphs[0]
    s0.text = "Model Interpretation"
    s0.font.bold = True
    s0.font.size = Pt(14)

    s1 = summary_tf.add_paragraph()
    s1.text = summary_text
    s1.level = 0  # <-- no bullet
    s1.font.size = Pt(14)


def _add_cadet_afsc_utility_chart_slide(instance, cadet, prs):

    slide = prs.slides.add_slide(prs.slide_layouts[3])  # chart_slide_layout
    title = slide.shapes.title
    title.text = f"Cadet '{cadet}' Calculated AFSC Utility"

    # Add the image to the slide
    for shape in slide.placeholders:
        if "Picture" in shape.name:
            # Create in-memory binary stream
            image_stream = io.BytesIO()

            # Save figure into memory (NOT disk)
            chart_obj = afccp.visualizations.charts.CadetUtilityGraph(instance, cadet=cadet, overwrite_save=True)
            chart_obj.fig.savefig(image_stream)  # , format='png', bbox_inches='tight', dpi=300)
            image_stream.seek(0)

            # Add to slide
            shape.insert_picture(image_stream)

            # Close figure to avoid memory leaks
            plt.close(chart_obj.fig)


def _add_cadet_state_utility_chart_slide(instance, cadet, prs):

    slide = prs.slides.add_slide(prs.slide_layouts[3])  # chart_slide_layout
    title = slide.shapes.title
    title.text = f"Cadet '{cadet}' Utility by State"

    # Add the image to the slide
    for shape in slide.placeholders:
        if "Picture" in shape.name:
            # Create in-memory binary stream
            image_stream = io.BytesIO()

            # Save figure into memory (NOT disk). # afccp.visualizations.charts.
            chart_obj = afccp.visualizations.charts.CadetStateUtilityGraph(instance, cadet=cadet, overwrite_save=True)
            chart_obj.fig.savefig(image_stream)  # , format='png', bbox_inches='tight', dpi=300)
            image_stream.seek(0)

            # Add to slide
            shape.insert_picture(image_stream)

            # Close figure to avoid memory leaks
            plt.close(chart_obj.fig)


def _add_base_utility_chart_slide(instance, cadet, prs):

    slide = prs.slides.add_slide(prs.slide_layouts[3])  # chart_slide_layout
    title = slide.shapes.title
    title.text = f"Cadet '{cadet}' Base Utility"

    # Add the image to the slide
    for shape in slide.placeholders:
        if "Picture" in shape.name:
            # Create in-memory binary stream
            image_stream = io.BytesIO()

            # Save figure into memory (NOT disk). # afccp.visualizations.charts.
            chart_obj = afccp.visualizations.charts.CadetBaseUtilityGraph(instance, cadet=cadet, overwrite_save=True)
            chart_obj.fig.savefig(image_stream)  # , format='png', bbox_inches='tight', dpi=300)
            image_stream.seek(0)

            # Add to slide
            shape.insert_picture(image_stream)

            # Close figure to avoid memory leaks
            plt.close(chart_obj.fig)


def _add_course_utility_chart_slide(instance, cadet, prs):

    slide = prs.slides.add_slide(prs.slide_layouts[3])  # chart_slide_layout
    title = slide.shapes.title
    title.text = f"Cadet '{cadet}' Course Utility"

    # Add the image to the slide
    for shape in slide.placeholders:
        if "Picture" in shape.name:
            # Create in-memory binary stream
            image_stream = io.BytesIO()

            # Save figure into memory (NOT disk). # afccp.visualizations.charts.
            chart_obj = afccp.visualizations.charts.CadetCourseUtilityGraph(instance, cadet=cadet, overwrite_save=True)
            chart_obj.fig.savefig(image_stream)  # , format='png', bbox_inches='tight', dpi=300)
            image_stream.seek(0)

            # Add to slide
            shape.insert_picture(image_stream)

            # Close figure to avoid memory leaks
            plt.close(chart_obj.fig)


def _generate_cadet_summary_text(p, i):
    # ----------------------------------------------------------
    # Core Data
    # ----------------------------------------------------------

    selected = p['J^Selected'][i]
    eligible = p['J^E'][i]
    eligible_selected = np.intersect1d(selected, eligible)

    selected_count = len(selected)
    eligible_count = len(eligible_selected)
    eligible_pct = round(100 * eligible_count / selected_count, 1) if selected_count > 0 else 0

    wA = p['weight_afsc'][i]
    wB = p['weight_base'][i]
    wC = p['weight_course'][i]

    bt = p['base_threshold'][i]
    tt = p['training_threshold'][i]
    training_pref = p['training_preferences'][i]

    # ----------------------------------------------------------
    # Determine Activation Timing
    # ----------------------------------------------------------

    def threshold_stage(threshold):
        if threshold >= 90:
            return "late"
        elif threshold >= 60:
            return "mid"
        elif threshold > 0:
            return "early"
        else:
            return "immediate"

    base_stage = threshold_stage(bt)
    train_stage = threshold_stage(tt)

    # ----------------------------------------------------------
    # Determine Early-Stage Dominance
    # ----------------------------------------------------------

    if base_stage in ["late"] and train_stage in ["late"]:
        early_structure = "AFSC-driven in early decision tiers"
    elif base_stage in ["early", "immediate"] and wB > wA:
        early_structure = "geographically sensitive early in the ranking structure"
    elif train_stage in ["early", "immediate"] and wC > wA:
        early_structure = "training-timeline sensitive early in the ranking structure"
    else:
        early_structure = "primarily AFSC-driven in early tiers with secondary objectives activating later"

    # ----------------------------------------------------------
    # Overall Weight Emphasis (Adjusted by Threshold)
    # ----------------------------------------------------------

    if wA >= max(wB, wC) and base_stage == "late" and train_stage == "late":
        priority_statement = "exhibits strong AFSC preference dominance"
    elif wB > wA and base_stage in ["early", "mid"]:
        priority_statement = "places meaningful emphasis on geographic outcomes"
    elif wC > wA and train_stage in ["early", "mid"]:
        priority_statement = "is meaningfully influenced by training timeline considerations"
    elif wB > wA and base_stage == "late":
        priority_statement = "indicates geographic importance, but only after AFSC thresholds are satisfied"
    elif wC > wA and train_stage == "late":
        priority_statement = "values training timing, though only in lower-ranked AFSC outcomes"
    else:
        priority_statement = "balances AFSC, base, and training objectives in a layered structure"

    # ----------------------------------------------------------
    # Threshold Interaction Insight
    # ----------------------------------------------------------

    if bt >= 90 and tt >= 90:
        threshold_statement = (
            "Base and training objectives activate only after high AFSC utility satisfaction, "
            "reinforcing a top-heavy AFSC preference structure."
        )
    elif bt < tt:
        threshold_statement = (
            "Geographic preferences activate earlier than training objectives, "
            "introducing location tradeoffs sooner in the ranking progression."
        )
    elif tt < bt:
        threshold_statement = (
            "Training timeline considerations activate before geographic preferences, "
            "influencing mid-tier AFSC allocations."
        )
    else:
        threshold_statement = (
            "Base and training thresholds activate at similar utility levels, "
            "producing a coordinated tradeoff structure."
        )

    # ----------------------------------------------------------
    # Training Timing Preference
    # ----------------------------------------------------------

    if training_pref == "Early":
        timing_statement = "The cadet prefers earlier training start dates."
    elif training_pref == "Late":
        timing_statement = "The cadet prefers later training start dates."
    else:
        timing_statement = "The cadet expressed neutral training start preferences."

    # ----------------------------------------------------------
    # Eligibility Pressure
    # ----------------------------------------------------------

    if eligible_pct == 100:
        eligibility_statement = "All submitted AFSC selections are eligible, minimizing constraint pressure."
    elif eligible_pct >= 75:
        eligibility_statement = "Most submitted AFSC selections are eligible, though minor constraint pressure exists."
    elif eligible_pct >= 50:
        eligibility_statement = "Eligibility constraints may materially influence allocation outcomes."
    else:
        eligibility_statement = "Significant eligibility constraints are likely to shape the final assignment."

    # ----------------------------------------------------------
    # Final Executive Summary
    # ----------------------------------------------------------

    summary = (
        f"This cadet submitted {selected_count} AFSC selections, "
        f"of which {eligible_count} ({eligible_pct}%) are eligible. "
        f"The preference structure is {early_structure}. "
        f"Overall, the cadet {priority_statement}. "
        f"{threshold_statement} "
        f"{timing_statement} "
        f"{eligibility_statement}"
    )

    return summary

