from pptx import Presentation
from pptx.util import Inches
import os
import copy
import numpy as np
import pandas as pd

# afccp modules
import afccp.core.globals


def generate_results_slides(instance):
    """
    Function to generate the results slides for a particular problem instance with solution
    """

    # Shorthand
    mdl_p = instance.mdl_p

    # Build the presentation object
    prs = Presentation(afccp.core.globals.paths['files'] + 'results_slide_template.pptx')

    # Delete the current slides in the template
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    # Slide Layouts
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]
    chart_slide_layout = prs.slide_layouts[2]
    bubble_slide_layout = prs.slide_layouts[3]
    closing_slide_layout = prs.slide_layouts[4]

    # Set width and height of presentation
    prs.slide_width = Inches(mdl_p['b_figsize'][0])
    prs.slide_height = Inches(mdl_p['b_figsize'][1])

    # Add title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = instance.data_name + " Classification Results (" + instance.solution['name'] + ")"
    # print(len(slide.placeholders))
    # print(slide.placeholders)
    # content = slide.placeholders[1]
    # content.text = "Rank, Name\nAFPC/DSYA\nTBD"

    # Initial content slide
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    title.text = "Overview"
    content = slide.placeholders[1]
    content.text = "Here's where the overview goes"

    # # Size of the chart
    # top, left = Inches(instance.mdl_p['ch_top']), Inches(instance.mdl_p['ch_left'])
    # height, width = Inches(instance.mdl_p['ch_height']), Inches(instance.mdl_p['ch_width'])

    # Size of the chart
    top, left = Inches(0), Inches(0)
    height, width = Inches(mdl_p['figsize'][1]), Inches(mdl_p['figsize'][0])

    # Get the file paths to all the relevant images
    folder_path = instance.export_paths['Analysis & Results'] + instance.solution_name + "/"
    folder = os.listdir(folder_path)
    chart_paths = {}
    for file in folder:
        if instance.solution_name in file and '.png' in file:
            chart_paths[file] = folder_path + file

    # Loop through the pictures I want in the order I want
    chart_text_order = {"PGL": ["Combined Quota", "quantity_bar"], "BUBBLE CHART": ["Cadet Choice.png"],
                        "SOC by Accessions Group": ["Accessions", "SOC Chart"],
                        "Gender by Accessions Group": ["Accessions", "Gender Chart"],
                        "Race by Accessions Group": ["Accessions", "Race Chart"],
                        "Ethnicity by Accessions Group": ["Accessions", "Ethnicity Chart"],
                        "SOC by AFSC": ["Extra Measure", "SOC Chart_proportion"],
                        "Gender by AFSC": ["Extra Measure", "Gender Chart_proportion"],
                        "Race by AFSC": ["Extra Measure", "Race Chart_proportion"],
                        "Ethnicity by AFSC": ["Extra Measure", "Ethnicity Chart_proportion"],
                        "Cadet Preference by AFSC": ["Utility", "quantity_bar_choice"]}
    for title_text in chart_text_order:

        # Loop through each potential image
        for file in chart_paths:

            # Determine if this is the image I want here
            found = True
            for test_text in chart_text_order[title_text]:
                if test_text not in file:
                    found = False
                    break

            # If this is the file I want, we do things
            if found:

                # Determine the layout of the chart needed
                if title_text == "BUBBLE CHART":

                    # Create the bubble chart slide
                    slide = prs.slides.add_slide(bubble_slide_layout)

                    # Add the image to the slide
                    for shape in slide.placeholders:
                        if "Picture" in shape.name:
                            shape.insert_picture(chart_paths[file])
                else:

                    # Create the AFSC/Accessions Chart slide
                    slide = prs.slides.add_slide(chart_slide_layout)
                    title = slide.shapes.title
                    title.text = title_text

                    # Add the picture to the slide
                    slide.shapes.add_picture(chart_paths[file], left, top, height=height, width=width)

                # Break out of this loop since we found the image we want
                break


    # Loop through each image file path to add the image to the presentation
    for file in chart_paths:

        # Add an empty slide
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add the picture to the slide
        slide.shapes.add_picture(chart_paths[file], left, top, height=height, width=width)

    # Add closing slide
    prs.slides.add_slide(closing_slide_layout)

    # Save the PowerPoint
    filepath = instance.export_paths['Analysis & Results'] + instance.solution_name + '/' + \
               instance.data_name + ' ' + instance.solution_name + '.pptx'
    prs.save(filepath)

def generate_comparison_slides(instance):
    """
    Function to generate the results slides for a particular problem instance with solution
    """

    # Shorthand
    mdl_p = instance.mdl_p

    # Build the presentation object
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    # Set width and height of presentation
    prs.slide_width = Inches(mdl_p['b_figsize'][0])
    prs.slide_height = Inches(mdl_p['b_figsize'][1])

    # Size of the chart
    top, left = Inches(0), Inches(0)

    # Get the file paths to all the relevant images
    folder_path = instance.export_paths['Analysis & Results'] + "Comparison Charts/"
    folder = os.listdir(folder_path)

    # Sort the files in the folder according to my preferred method
    chart_paths = {}
    keyword_order = [['Combined Quota'], ['Tier 1'], ['Male'], ['USAFA Proportion'], ['Merit'], ['Norm Score'],
                     ['Utility', 'dot'], ['Utility', 'mean_preference'], ['Utility', 'median_preference'],
                     ['Utility', 'Histogram'], ['Pareto', 'Cadets.png'], ['Pareto', ').png'], ['Similarity']]

    # Loop through each set of "keywords"
    for keywords in keyword_order:

        # Loop through each file until we have a match
        for file in folder:

            # Loop through all keywords to see if we have a match
            match = True
            for keyword in keywords:
                if keyword not in file:
                    match = False
                    break

            # If we have a "match", add it!
            if match:
                chart_paths[file] = folder_path + file
                folder.remove(file)  # Remove this file since we've accounted for it
                break

    # Add any remaining files
    for file in folder:
        if '.png' in file:
            chart_paths[file] = folder_path + file

    # Loop through each image file path to add the image to the presentation
    for file in chart_paths:

        # Add an empty slide
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add the picture to the slide
        slide.shapes.add_picture(chart_paths[file], left, top)  #, height=height, width=width)

    # Save the PowerPoint
    filepath = instance.export_paths['Analysis & Results'] + instance.data_name + ' ' + 'Comparison.pptx'
    prs.save(filepath)

def create_animation_slides(instance):
    """
    Function to generate the results slides for a particular problem instance with solution
    """

    # Shorthand
    mdl_p, sequence = instance.mdl_p, instance.solution['iterations']['sequence']

    # Build the presentation object
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    # Set width and height of presentation
    prs.slide_width = Inches(mdl_p['b_figsize'][0])
    prs.slide_height = Inches(mdl_p['b_figsize'][1])

    # Size of the chart
    top, left = Inches(0), Inches(0)
    height, width = Inches(mdl_p['b_figsize'][1]), Inches(mdl_p['b_figsize'][0])

    # Make sure we have the "sequence" folder
    if sequence not in os.listdir(instance.export_paths['Analysis & Results'] + "Cadet Board/"):
        raise ValueError("Error. Sequence folder '" + sequence + "' not in 'Cadet Board' folder.")

    # Get the file path to the sequence folder
    sequence_folder_path = instance.export_paths['Analysis & Results'] + "Cadet Board/" + sequence + '/'

    # Make sure we have the "focus" folder
    if mdl_p['focus'] not in os.listdir(sequence_folder_path):
        raise ValueError("Error. Focus folder '" + mdl_p['focus'] + "' not in '" + sequence + "' sequence folder.")

    # Files in the focus folder
    focus_folder_path = sequence_folder_path + mdl_p['focus'] + '/'
    folder = np.array(os.listdir(focus_folder_path))

    # Regular solution iterations frames: 1.png for example
    if "1.png" in folder:

        # Sort the folder in order by the frames
        int_vals = np.array([int(file[0]) for file in folder])
        indices = np.argsort(int_vals)
        img_paths = {file: focus_folder_path + file for file in folder[indices]}

    # Matching Algorithm proposals/rejections frames: 1 (Proposals).png & 1 (Rejections).png for example
    else:

        # The integer values at the beginning of each file
        int_vals = np.array([int(file[:2]) for file in folder])
        min, max = np.min(int_vals), np.max(int_vals)

        # Loop through the files to get the ordered list of frames
        img_paths = {}
        for val in np.arange(min, max + 1):
            indices = np.where(int_vals == val)[0]
            files = folder[indices]
            if len(files) > 1:
                if 'Proposals' in files[0]:
                    img_paths[files[0]] = focus_folder_path + files[0]
                    img_paths[files[1]] = focus_folder_path + files[1]
                else:
                    img_paths[files[1]] = focus_folder_path + files[1]
                    img_paths[files[0]] = focus_folder_path + files[0]
            else:
                img_paths[files[0]] = focus_folder_path + files[0]

    # Loop through each image file path to add the image to the presentation
    for file in img_paths:

        # Add an empty slide
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add the picture to the slide
        slide.shapes.add_picture(img_paths[file], left, top, height=height, width=width)

    # Save the PowerPoint
    filepath = sequence_folder_path + mdl_p['focus'] + '.pptx'
    prs.save(filepath)